import os
import io
import json
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from google.generativeai.types import generation_types
from pydantic import BaseModel
from pptx import Presentation

app = FastAPI(title="Papipoint API", version="0.1.0")

# Allow CORS for Next.js frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")

class Slide(BaseModel):
    title: str
    subtitle: str | None = None
    bullets: list[str] = []

class PresentationStructure(BaseModel):
    title: str
    slides: list[Slide]

@app.get("/")
def read_root():
    return {"message": "Welcome to the Papipoint API"}

@app.post("/api/process-pptx")
async def process_pptx(
    file: UploadFile = File(...),
    custom_api_key: str | None = Form(None)
):
    """
    Endpoint to receive a .pptx file, extract text, and use AI to restructure it.
    """
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported.")
    
    # 1. Extract Text
    content = await file.read()
    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid PPTX file: {str(e)}")

    extracted_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {i+1} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        extracted_text.append("\n".join(slide_text))
    
    full_text = "\n\n".join(extracted_text)
    
    if not full_text.strip():
        raise HTTPException(status_code=400, detail="No text could be extracted from the presentation.")

    # 2. Call AI
    api_key_to_use = custom_api_key if custom_api_key else GEMINI_API_KEY
    if not api_key_to_use:
        raise HTTPException(status_code=500, detail="No API key available for AI generation.")
    
    genai.configure(api_key=api_key_to_use)
    
    # We use gemini-1.5-flash as it is fast and supports JSON response schema well.
    model = genai.GenerativeModel('gemini-1.5-flash')
    
    prompt = f"""
    You are an expert presentation designer. I have extracted raw, unformatted text from an old PowerPoint presentation.
    Your job is to read this text, understand the core narrative, and restructure it into a modern, clean presentation outline.
    
    Remove fluff. Keep bullet points concise (max 3-5 per slide). Group related concepts.
    Output strictly in JSON matching the schema provided.

    JSON Schema Required:
    {
      "title": "Presentation Main Title",
      "slides": [
        {
          "title": "Slide Title",
          "subtitle": "Optional slide subtitle",
          "bullets": ["Bullet 1", "Bullet 2"]
        }
      ]
    }

    Raw Extracted Text:
    {full_text}
    """

    try:
        response = model.generate_content(
            prompt,
            generation_config=genai.GenerationConfig(
                response_mime_type="application/json",
            )
        )
        
        # Parse the JSON response
        result_json = json.loads(response.text)
        return result_json
        
    except Exception as e:
        print("AI Generation Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to process with AI: {str(e)}")
